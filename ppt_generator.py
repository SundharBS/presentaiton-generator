from pptx import Presentation


def safe_replace(text, replacements):

    if not text:
        return text

    for key, value in replacements.items():

        text = text.replace(
            key,
            str(value)
        )

    return text


def replace_text_in_shape(shape, replacements):

    if not hasattr(shape, "text_frame"):
        return

    text_frame = shape.text_frame

    for paragraph in text_frame.paragraphs:

        for run in paragraph.runs:

            original_text = run.text

            updated_text = safe_replace(
                original_text,
                replacements
            )

            if updated_text != original_text:

                run.text = updated_text


def generate_ppt(
    company_name,
    analysis,
    selected_sections
):

    # LOAD TEMPLATE
    prs = Presentation(
        "template.pptx"
    )

    replacements = {

        # COVER PAGE
        "{{COMPANY_NAME}}":
        company_name,

        "{{COMPANY_TAGLINE}}":
        "AI Generated Investment Banking Deck",

        "{{PRESENTATION_TITLE}}":
        "Strategic Overview & Investment Thesis",

        "{{STAT1_VALUE}}":
        analysis.get(
            "aum",
            "N/A"
        ),

        "{{STAT1_LABEL}}":
        "Revenue / AUM",

        "{{STAT2_VALUE}}":
        analysis.get(
            "growth",
            "N/A"
        ),

        "{{STAT2_LABEL}}":
        "Growth",

        "{{STAT3_VALUE}}":
        analysis.get(
            "gnpa",
            "N/A"
        ),

        "{{STAT3_LABEL}}":
        "Risk Indicator",

        "{{STAT4_VALUE}}":
        analysis.get(
            "capital",
            "N/A"
        ),

        "{{STAT4_LABEL}}":
        "Capital Base",

        "{{IB_DIVISION_NAME}}":
        "Investment Banking Division",

        "{{PRESENTATION_DATE}}":
        "May 2026",

        # TOC
        "{{CLIENT_SITUATIONAL_ANALYSIS}}":
        "Client Situational Analysis",

        "{{CLIENT_SITUATIONAL_ANALYSIS_SUBTITLE}}":
        "Company Overview & Current Position",

        "{{SWOT_ANALYSIS}}":
        "SWOT Analysis",

        "{{SWOT_ANALYSIS_SUBTITLE}}":
        "Strengths, Weaknesses, Opportunities & Threats",

        "{{FINANCIAL_PERFORMANCE_OVERVIEW}}":
        "Financial Performance Overview",

        "{{FINANCIAL_PERFORMANCE_OVERVIEW_SUBTITLE}}":
        "Growth & Operating Metrics",

        "{{COMPETITIVE_POSITION}}":
        "Competitive Position",

        "{{COMPETITIVE_POSITION_SUBTITLE}}":
        "Positioning & Competitive Moat",

        "{{MARKET_INDUSTRY_OVERVIEW}}":
        "Market & Industry Overview",

        "{{MARKET_INDUSTRY_OVERVIEW_SUBTITLE}}":
        "Industry Landscape & Opportunity",

        "{{KEY_CONSIDERATIONS_RISK_FACTORS}}":
        "Key Considerations & Risk Factors",

        "{{KEY_CONSIDERATIONS_RISK_FACTORS_SUBTITLE}}":
        "Key Risks & Mitigants",

        "{{STRATEGIC_OPTIONS_THESIS}}":
        "Strategic Options / Thesis",

        "{{STRATEGIC_OPTIONS_THESIS_SUBTITLE}}":
        "Future Roadmap & Strategic Direction",

        "{{VALUATION_ANALYSIS}}":
        "Valuation Analysis",

        "{{VALUATION_ANALYSIS_SUBTITLE}}":
        "Peer Benchmarking & Investment View",

        # SECTIONS
        "{{SECTION_01_HEADER}}":
        "01 | CLIENT SITUATIONAL ANALYSIS",

        "{{SECTION_01_TITLE}}":
        analysis.get(
            "client_situational_analysis",
            "Not Available"
        ),

        "{{SECTION_02_HEADER}}":
        "02 | SWOT ANALYSIS",

        "{{SECTION_03_HEADER}}":
        "03 | FINANCIAL PERFORMANCE",

        "{{SECTION_04_HEADER}}":
        "04 | COMPETITIVE POSITION",

        "{{SECTION_05_HEADER}}":
        "05 | MARKET OPPORTUNITY",

        "{{SECTION_06_HEADER}}":
        "06 | RISK FACTORS & MITIGANTS",

        "{{SECTION_07_HEADER}}":
        "07 | STRATEGIC OPTIONS / THESIS",

        "{{SECTION_08_HEADER}}":
        "08 | VALUATION ANALYSIS",

        # BODY CONTENT
        "{{MARKET_INDUSTRY_OVERVIEW}}":
        analysis.get(
            "market_industry_overview",
            "Not Available"
        ),

        "{{STRATEGIC_OPTIONS_THESIS}}":
        analysis.get(
            "strategic_options_thesis",
            "Not Available"
        ),

        "{{VALUATION_ANALYSIS}}":
        analysis.get(
            "valuation_analysis",
            "Not Available"
        ),

        "{{KEY_CONSIDERATIONS_RISK_FACTORS}}":
        analysis.get(
            "key_considerations_risk_factors",
            "Not Available"
        ),

        "{{COMPANY_TAGLINE_CLOSING}}":
        "Expanding Horizons",

        "{{REPORT_YEAR}}":
        "FY2025"
    }

    # REPLACE TEXT WHILE PRESERVING FORMAT
    for slide in prs.slides:

        for shape in slide.shapes:

            replace_text_in_shape(
                shape,
                replacements
            )

    ppt_path = (
        f"{company_name}_Investment_Deck.pptx"
    )

    prs.save(ppt_path)

    return ppt_path